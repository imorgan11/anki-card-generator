import streamlit as st
import pandas as pd
import json
import genanki
import random
import os
import re
from google import genai
from google.genai import types
from docx import Document
from pypdf import PdfReader
from pptx import Presentation


# ------------------------------------------------------------------------------
# Helper functions
# ------------------------------------------------------------------------------

def clean_file_name(name):
    """
    Makes the deck name safe to use as a file name.
    """
    cleaned_name = re.sub(r'[^a-zA-Z0-9_-]', '_', name)
    return cleaned_name


def clean_ai_json(ai_output):
    """
    Cleans Gemini output and extracts the JSON list from the response.
    """
    ai_output = ai_output.strip()

    # Remove markdown code fences if present
    if ai_output.startswith('```json'):
        ai_output = ai_output.replace('```json', '').replace('```', '').strip()

    elif ai_output.startswith('```'):
        ai_output = ai_output.replace('```', '').strip()

    # Extract only the JSON array if Gemini adds extra text before or after it
    start_index = ai_output.find('[')
    end_index = ai_output.rfind(']')

    if start_index != -1 and end_index != -1:
        ai_output = ai_output[start_index:end_index + 1]

    return ai_output


def create_anki_package(cards_df, deck_name):
    """
    Creates a real .apkg Anki deck file from a DataFrame with Question and Answer columns.
    """

    model_id = random.randrange(1 << 30, 1 << 31)
    deck_id = random.randrange(1 << 30, 1 << 31)

    anki_model = genanki.Model(
        model_id,
        'Basic Question Answer Model',
        fields=[
            {'name': 'Question'},
            {'name': 'Answer'}
        ],
        templates=[
            {
                'name': 'Card 1',
                'qfmt': '{{Question}}',
                'afmt': '{{FrontSide}}<hr id="answer">{{Answer}}'
            }
        ]
    )

    anki_deck = genanki.Deck(
        deck_id,
        deck_name
    )

    for index, row in cards_df.iterrows():
        question = str(row['Question'])
        answer = str(row['Answer'])

        note = genanki.Note(
            model=anki_model,
            fields=[question, answer]
        )

        anki_deck.add_note(note)

    safe_deck_name = clean_file_name(deck_name)
    output_file = f'{safe_deck_name}_anki_deck.apkg'

    package = genanki.Package(anki_deck)
    package.write_to_file(output_file)

    return output_file


def extract_text_from_txt(file):
    """
    Extracts text from a .txt file.
    """
    return file.read().decode('utf-8')


def extract_text_from_docx(file):
    """
    Extracts typed paragraph text from a Word document.
    Also counts images.
    """
    document = Document(file)

    paragraphs = []

    for paragraph in document.paragraphs:
        paragraphs.append(paragraph.text)

    extracted_text = '\n'.join(paragraphs)

    image_count = 0

    for relationship in document.part.rels.values():
        if 'image' in relationship.reltype:
            image_count += 1

    return extracted_text, image_count


def extract_text_from_pdf(file):
    """
    Extracts text from a PDF.
    Does not OCR scanned/image-based PDFs.
    """
    reader = PdfReader(file)

    pages = []

    for page in reader.pages:
        page_text = page.extract_text()

        if page_text is not None:
            pages.append(page_text)

    extracted_text = '\n'.join(pages)

    return extracted_text


def extract_text_from_pptx(file):
    """
    Extracts typed text from a PowerPoint.
    Also counts images.
    """
    presentation = Presentation(file)

    slide_texts = []
    image_count = 0

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_texts.append(f'--- Slide {slide_number} ---')

        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                if shape.text.strip() != '':
                    slide_texts.append(shape.text)

            if shape.shape_type == 13:
                image_count += 1

    extracted_text = '\n'.join(slide_texts)

    return extracted_text, image_count


# ------------------------------------------------------------------------------
# Gemini API key
# ------------------------------------------------------------------------------

gemini_api_key = st.secrets["GEMINI_API_KEY"]


# ------------------------------------------------------------------------------
# Welcome page
# ------------------------------------------------------------------------------

st.title('Welcome to Isaac\'s Anki Card Generator!')

st.write('This is a beta version designed to help friends only.')

st.write(
    'This app will accept uploads in .txt, .pdf, .docx, and .pptx formats, '
    'and then create flashcards of your choice based on the material. '
    'Flashcards can be previewed and then downloaded as an Anki deck package. Enjoy!'
)


# ------------------------------------------------------------------------------
# Initializing deck name
# ------------------------------------------------------------------------------

deck_name = st.text_input('Type deck name here:')

if deck_name.strip() != '':
    st.write(f'Deck name: {deck_name}')


# ------------------------------------------------------------------------------
# Card type selection
# ------------------------------------------------------------------------------

card_type = st.selectbox(
    'Select card type:',
    [
        'Definition',
        'Application',
        'Distinction',
        'Mixed'
    ]
)

st.write(f'Selected card type: {card_type}')


# ------------------------------------------------------------------------------
# Material upload and text extraction
# ------------------------------------------------------------------------------

final_study_text = ''

file = st.file_uploader(
    'Upload notes, slides, and anything else here:',
    type=['txt', 'pdf', 'docx', 'pptx']
)

if file is not None:
    st.write(f'File name: {file.name}')
    st.write(f'File type: {file.type}')
    st.write(f'File size: {file.size} bytes')

    extracted_text = ''

    # --------------------------------------------------------------------------
    # TXT extraction
    # --------------------------------------------------------------------------

    if file.name.endswith('.txt'):
        extracted_text = extract_text_from_txt(file)

    # --------------------------------------------------------------------------
    # DOCX extraction
    # --------------------------------------------------------------------------

    elif file.name.endswith('.docx'):
        extracted_text, image_count = extract_text_from_docx(file)

        if image_count > 0:
            st.warning(
                f'This Word document contains {image_count} image(s). '
                'Text inside images will not be extracted yet.'
            )

    # --------------------------------------------------------------------------
    # PDF extraction
    # --------------------------------------------------------------------------

    elif file.name.endswith('.pdf'):
        extracted_text = extract_text_from_pdf(file)

        if extracted_text.strip() == '':
            st.warning(
                'No readable text was extracted from this PDF. '
                'It may be a scanned or image-based PDF.'
            )

    # --------------------------------------------------------------------------
    # PPTX extraction
    # --------------------------------------------------------------------------

    elif file.name.endswith('.pptx'):
        extracted_text, image_count = extract_text_from_pptx(file)

        if image_count > 0:
            st.warning(
                f'This PowerPoint contains {image_count} image(s). '
                'Text inside images will not be extracted yet.'
            )

    # --------------------------------------------------------------------------
    # Unsupported extraction
    # --------------------------------------------------------------------------

    else:
        st.warning(
            'Text extraction is only set up for .txt, .docx, .pdf, and .pptx files right now.'
        )

    # --------------------------------------------------------------------------
    # Editable extracted text
    # --------------------------------------------------------------------------

    if extracted_text.strip() != '':
        st.subheader('Extracted Text Preview')

        final_study_text = st.text_area(
            'Extracted text',
            extracted_text,
            height=300
        )


# ------------------------------------------------------------------------------
# Final study text check
# ------------------------------------------------------------------------------

if file is not None:
    st.write(f'Final study text length: {len(final_study_text)} characters')


# ------------------------------------------------------------------------------
# Generate Anki cards
# ------------------------------------------------------------------------------

if st.button('Click here to generate cards!'):

    if gemini_api_key.strip() == '':
        st.warning('Please enter your Gemini API key.')

    elif deck_name.strip() == '':
        st.warning('Please enter a deck name.')

    elif final_study_text.strip() == '':
        st.warning('Please upload a file and make sure text was extracted.')

    else:
        client = genai.Client(api_key=gemini_api_key)

        prompt = f"""
You are an expert Anki flashcard generator.

Create high-quality Anki flashcards from the study material below.

Card type selected by user: {card_type}

Rules:
- Return only valid JSON.
- Do not use markdown.
- Do not include explanations outside the JSON.
- The JSON must be a list of objects.
- Each object must have exactly two keys: "Question" and "Answer".
- Questions should be clear and testable.
- Answers should be concise but complete.
- Avoid duplicate cards.
- Make enough cards to cover the important material.
- If the selected card type is Definition, prioritize definition cards.
- If the selected card type is Application, prioritize applied scenario cards.
- If the selected card type is Distinction, prioritize comparison cards.
- If the selected card type is Mixed, include definition, application, and distinction cards.
- Do not create cards from irrelevant headers, page numbers, or formatting artifacts.
- Keep answers brief enough for flashcard use, but not so brief that they lose key meaning.
- Your entire response must start with [ and end with ].
- Do not write "Here are the cards".
- Do not include notes, comments, or explanations.
- Escape all quotation marks inside answers properly.

Study material:
{final_study_text}
"""

        try:
            with st.spinner('Generating cards...'):
                response = client.models.generate_content(
                    model='gemini-2.5-flash',
                    contents=prompt,
                    config=types.GenerateContentConfig(
                    response_mime_type='application/json'
                    )
                )

            ai_output = response.text.strip()
            ai_output = clean_ai_json(ai_output)

            cards = json.loads(ai_output)

            cards_df = pd.DataFrame(cards)

            # Make sure required columns exist
            if 'Question' not in cards_df.columns or 'Answer' not in cards_df.columns:
                st.error('Gemini returned JSON, but it did not contain Question and Answer columns.')
                st.write(cards_df)

            else:
                cards_df = cards_df[['Question', 'Answer']]

                st.subheader('Generated Card Preview')
                st.dataframe(cards_df)

                # ------------------------------------------------------------------
                # APKG download
                # ------------------------------------------------------------------
                
                safe_deck_name = clean_file_name(deck_name)
                
                anki_package_file = create_anki_package(cards_df, deck_name)
                
                with open(anki_package_file, 'rb') as apkg_file:
                    st.download_button(
                        label='Download Anki Deck Package',
                        data=apkg_file,
                        file_name=f'{safe_deck_name}_anki_deck.apkg',
                        mime='application/octet-stream'
                    )

        except json.JSONDecodeError:
            st.error('Gemini did not return valid JSON. Raw response shown below:')
            st.write(ai_output)

        except Exception as error:
            st.error('Something went wrong while generating cards.')
            st.write(error)
